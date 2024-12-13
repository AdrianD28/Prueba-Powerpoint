# tests/test_main.py
import os
from pptx import Presentation
from app.main import crear_pptx, modificar_pptx

def test_crear_y_modificar_pptx():
    # Ejecutar las funciones
    crear_pptx()
    modificar_pptx()
    
    # Verificar que el archivo se ha creado
    assert os.path.exists('trial.pptx'), "El archivo trial.pptx no se creó."

    # Verificar que el texto se ha modificado
    prs = Presentation('trial.pptx')
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    assert run.text == "Gracias", "El texto no se modificó correctamente."

# Ejecutar la prueba
if __name__ == "__main__":
    test_crear_y_modificar_pptx()
    print("Todas las pruebas pasaron correctamente.")