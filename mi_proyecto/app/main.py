from pptx import Presentation
from pptx.enum.text import PP_ALIGN

def crear_pptx():
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
    text_frame = textbox.text_frame
    text_frame.text = "Hola"
    prs.save('trial.pptx')

def modificar_pptx():
    prs = Presentation('trial.pptx')
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text == "Hola":
                        run.text = "Gracias"
                # Centrar el texto en el cuadro de texto
                paragraph.alignment = PP_ALIGN.CENTER
    prs.save('trial.pptx')

crear_pptx()
modificar_pptx()